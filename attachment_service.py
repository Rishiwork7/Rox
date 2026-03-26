import io
import os
import random
import shutil
import string
import tempfile
import subprocess
import sys
from pathlib import Path

from html2image import Html2Image
from PIL import Image

try:
    from weasyprint import HTML
except ImportError:
    HTML = None

try:
    import docx
    from docx.shared import Inches as DocxInches
except ImportError:
    docx = None
    DocxInches = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
except ImportError:
    Presentation = None
    PptxInches = None


DEFAULT_IMAGE_HEIGHT = 1080
SUPPORTED_IMAGE_FORMATS = {"png", "jpeg", "jpg", "tiff", "bmp", "gif", "webp"}
HTML2IMAGE_BROWSERS = ("chrome", "chromium", "edge")
HTML2IMAGE_FLAGS = [
    "--default-background-color=ffffffff",
    "--hide-scrollbars",
    "--no-sandbox",
    "--disable-gpu",
]
BROWSER_EXECUTABLES = {
    "chrome": (
        shutil.which("google-chrome"),
        shutil.which("chrome"),
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    ),
    "chromium": (
        shutil.which("chromium"),
        shutil.which("chromium-browser"),
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
    ),
    "edge": (
        shutil.which("microsoft-edge"),
        shutil.which("edge"),
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
    ),
}


def _random_invoice_id():
    return "INV-" + "".join(random.choices(string.ascii_uppercase + string.digits, k=6))


def _normalize_img_format(img_format):
    fmt = (img_format or "png").lower()
    return "jpeg" if fmt == "jpg" else fmt


def _resolve_image_size(config_dict):
    width = int(config_dict.get("width", 800) or 800)
    if config_dict.get("random_width"):
        width += random.randint(-50, 50)

    height = int(config_dict.get("height", 0) or 0)
    return max(width, 320), max(height, 1) if height > 0 else DEFAULT_IMAGE_HEIGHT


def _create_html2image(output_path, size):
    last_error = None
    for browser_name in HTML2IMAGE_BROWSERS:
        executable = next(
            (candidate for candidate in BROWSER_EXECUTABLES[browser_name] if candidate and Path(candidate).exists()),
            None,
        )
        try:
            hti = Html2Image(
                browser=browser_name,
                browser_executable=executable,
                output_path=str(output_path),
                size=size,
                keep_temp_files=False,
                custom_flags=HTML2IMAGE_FLAGS,
                disable_logging=True,
            )
            if hasattr(hti.browser, "use_new_headless"):
                hti.browser.use_new_headless = True
            return hti
        except Exception as exc:
            last_error = exc

    raise RuntimeError(
        "html2image could not initialize a supported browser (chrome/chromium/edge)"
    ) from last_error


def _render_html_to_base_png(final_attachment_html, config_dict):
    width, height = _resolve_image_size(config_dict)

    with tempfile.TemporaryDirectory(prefix="rox_attach_") as temp_dir:
        output_dir = Path(temp_dir)
        
        # Fallback: ensure /tmp/html2image exists if library insists on using it
        try:
            os.makedirs("/tmp/html2image", exist_ok=True)
        except:
            pass

        screenshot_name = "attachment_source.png"
        screenshot_path = output_dir / screenshot_name
        
        # Write HTML to a file in our controlled temp_dir
        html_file_path = output_dir / "attachment_source.html"
        with open(html_file_path, "w", encoding="utf-8") as f:
            f.write(final_attachment_html)

        hti = _create_html2image(output_dir, (width, height))

        generated_paths = hti.screenshot(
            html_file=str(html_file_path),
            save_as=screenshot_name,
            size=(width, height),
        )

        if generated_paths:
            screenshot_path = Path(generated_paths[0])

        if not screenshot_path.exists():
            raise RuntimeError("html2image did not generate the temporary PNG")

        with Image.open(screenshot_path) as source_image:
            base_image = source_image.copy()

        os.remove(screenshot_path)
        return base_image


def _render_html_to_pdf_direct(final_attachment_html):
    """
    Uses Headless Chrome/Chromium to print HTML directly to PDF.
    """
    executable = None
    for browser_name in HTML2IMAGE_BROWSERS:
        executable = next(
            (candidate for candidate in BROWSER_EXECUTABLES[browser_name] if candidate and Path(candidate).exists()),
            None,
        )
        if executable:
            break
            
    if not executable:
        raise RuntimeError("No browser found for direct PDF conversion")

    with tempfile.TemporaryDirectory(prefix="rox_pdf_") as temp_dir:
        output_dir = Path(temp_dir)
        html_file = output_dir / "input.html"
        pdf_file = output_dir / "output.pdf"
        
        with open(html_file, "w", encoding="utf-8") as f:
            f.write(final_attachment_html)
            
        args = [
            executable,
            "--headless",
            "--disable-gpu",
            "--no-sandbox",
            f"--print-to-pdf={pdf_file}",
            str(html_file)
        ]
        
        try:
            subprocess.run(args, check=True, capture_output=True, timeout=30)
            if pdf_file.exists():
                return pdf_file.read_bytes()
        except Exception as e:
            raise RuntimeError(f"Chrome PDF conversion failed: {str(e)}")
            
    raise RuntimeError("Failed to generate PDF via Chrome")


def _apply_image_options(image, target_format, config_dict):
    processed = image.copy()

    if config_dict.get("grayscale"):
        processed = processed.convert("L")

    if config_dict.get("crop"):
        width, height = processed.size
        processed = processed.crop((0, 0, width, int(height * 0.8)))

    # Ensure fallback to white background for transparency in most formats
    if processed.mode in {"RGBA", "LA", "P"}:
        # If it's a format that doesn't support transparency or if we just want to be safe
        # we flatten onto white.
        white_bg = Image.new("RGB", processed.size, "white")
        if processed.mode == "RGBA":
            white_bg.paste(processed, mask=processed.getchannel("A"))
        else:
            white_bg.paste(processed.convert("RGBA"), mask=processed.convert("RGBA").getchannel("A"))
        processed = white_bg

    return processed


def _image_to_bytes(image, target_format, config_dict):
    output = io.BytesIO()
    save_format = target_format.upper() if target_format != "jpeg" else "JPEG"
    save_kwargs = {"format": save_format}

    if target_format in {"jpeg", "webp"}:
        quality = 85
        if config_dict.get("hq"):
            quality = 95
        if config_dict.get("hsize"):
            quality = 100
        save_kwargs["quality"] = quality

    image.save(output, **save_kwargs)
    return output.getvalue()


def _render_html_to_image_bytes(final_attachment_html, img_format, config_dict):
    base_png = _render_html_to_base_png(final_attachment_html, config_dict)
    processed_image = _apply_image_options(base_png, img_format, config_dict)
    return _image_to_bytes(processed_image, img_format, config_dict)


def _image_bytes_to_pdf_bytes(image_bytes):
    with Image.open(io.BytesIO(image_bytes)) as image:
        pdf_ready_image = image.convert("RGB") if image.mode != "RGB" else image.copy()
        buffer = io.BytesIO()
        pdf_ready_image.save(buffer, format="PDF")
        return buffer.getvalue()


def _image_bytes_to_docx_bytes(image_bytes):
    if docx is None or DocxInches is None:
        raise RuntimeError("python-docx is not installed")

    document = docx.Document()
    document.add_picture(io.BytesIO(image_bytes), width=DocxInches(6))
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _image_bytes_to_pptx_bytes(image_bytes):
    if Presentation is None or PptxInches is None:
        raise RuntimeError("python-pptx is not installed")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(image_bytes), PptxInches(0.5), PptxInches(0.5), width=PptxInches(9))
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def generate_attachment(final_attachment_html, output_mode, img_format="png", config_dict=None):
    config_dict = config_dict or {}
    inv_no = config_dict.get("inv_no") or _random_invoice_id()
    normalized_format = _normalize_img_format(img_format)
    cid_id = None

    if output_mode == "Inline Html":
        return None, None, None

    if output_mode == "Raw pdf":
        try:
            pdf_bytes = _render_html_to_pdf_direct(final_attachment_html)
            return pdf_bytes, f"invoice_{inv_no}.pdf", None
        except Exception as e:
            # Fallback to WeasyPrint if available
            if HTML is not None:
                pdf_buffer = io.BytesIO()
                HTML(string=final_attachment_html).write_pdf(target=pdf_buffer)
                return pdf_buffer.getvalue(), f"invoice_{inv_no}.pdf", None
            # Last resort: Image-based PDF
            image_bytes = _render_html_to_image_bytes(final_attachment_html, "png", config_dict)
            return _image_bytes_to_pdf_bytes(image_bytes), f"invoice_{inv_no}.pdf", None

    if normalized_format not in SUPPORTED_IMAGE_FORMATS:
        raise ValueError(f"Unsupported image format: {img_format}")

    image_bytes = _render_html_to_image_bytes(final_attachment_html, normalized_format, config_dict)

    if output_mode == "To pdf":
        return _image_bytes_to_pdf_bytes(image_bytes), f"invoice_{inv_no}.pdf", None

    if output_mode == "To image":
        return image_bytes, f"invoice_{inv_no}.{normalized_format}", None

    if output_mode == "Inline image":
        cid_id = f"img_{inv_no.lower()}"
        return image_bytes, f"invoice_{inv_no}.{normalized_format}", cid_id

    if output_mode == "Docx":
        return _image_bytes_to_docx_bytes(image_bytes), f"invoice_{inv_no}.docx", None

    if output_mode == "PPTX":
        return _image_bytes_to_pptx_bytes(image_bytes), f"invoice_{inv_no}.pptx", None

    raise ValueError(f"Unsupported output mode: {output_mode}")


def generate_attachment_buffer(html_content, output_mode, img_format="png", config_dict=None):
    return generate_attachment(
        final_attachment_html=html_content,
        output_mode=output_mode,
        img_format=img_format,
        config_dict=config_dict,
    )
